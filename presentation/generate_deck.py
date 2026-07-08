from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- Theme ----------
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
SKY = RGBColor(0x2E, 0x86, 0xDE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x6B, 0x72, 0x80)
GREEN = RGBColor(0x1E, 0x8E, 0x3E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
# ---------- Helpers ----------
def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(BLANK)
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(11.7), Inches(1.5))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = WHITE
    box2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.7), Inches(1))
    p2 = box2.text_frame.paragraphs[0]; p2.text = subtitle
    p2.font.size = Pt(20); p2.font.color.rgb = SKY
    return slide

def add_section_header(title):
    slide = prs.slides.add_slide(BLANK)
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = SKY; bg.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.7), Inches(1.2))
    p = box.text_frame.paragraphs[0]; p.text = title
    p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = WHITE
    return slide

def _title_bar(slide, title):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.9))
    p = box.text_frame.paragraphs[0]; p.text = title
    p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = NAVY
    line = slide.shapes.add_shape(1, Inches(0.6), Inches(1.15), Inches(6.5), Pt(3))
    line.fill.solid(); line.fill.fore_color.rgb = SKY; line.line.fill.background()

def add_bullet_slide(title, bullets, note=None):
    slide = prs.slides.add_slide(BLANK)
    _title_bar(slide, title)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.3))
    tf = box.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        indent = b.startswith("  ")
        p.text = ("•  " if not indent else "‑  ") + b.strip()
        p.font.size = Pt(20) if not indent else Pt(17)
        p.font.color.rgb = NAVY if not indent else GREY
        p.space_after = Pt(12)
        p.level = 1 if indent else 0
    if note:
        nb = slide.shapes.add_textbox(Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.6))
        npar = nb.text_frame.paragraphs[0]; npar.text = note
        npar.font.size = Pt(13); npar.font.italic = True; npar.font.color.rgb = GREY
    return slide
def add_two_col_slide(title, left_heading, left_bullets, right_heading, right_bullets):
    slide = prs.slides.add_slide(BLANK)
    _title_bar(slide, title)
    def col(x, heading, bullets):
        hb = slide.shapes.add_textbox(Inches(x), Inches(1.6), Inches(5.6), Inches(0.5))
        hp = hb.text_frame.paragraphs[0]; hp.text = heading
        hp.font.size = Pt(20); hp.font.bold = True; hp.font.color.rgb = SKY
        bb = slide.shapes.add_textbox(Inches(x), Inches(2.15), Inches(5.6), Inches(4.8))
        tf = bb.text_frame; tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "•  " + b
            p.font.size = Pt(17); p.font.color.rgb = NAVY
            p.space_after = Pt(10)
    col(0.7, left_heading, left_bullets)
    col(6.9, right_heading, right_bullets)
    return slide

def add_status_slide(title, rows):
    # rows: list of (item, status, color) where color is GREEN/SKY/GREY
    slide = prs.slides.add_slide(BLANK)
    _title_bar(slide, title)
    y = 1.7
    for item, status, color in rows:
        box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(9.5), Inches(0.55))
        p = box.text_frame.paragraphs[0]; p.text = item
        p.font.size = Pt(18); p.font.color.rgb = NAVY
        chip = slide.shapes.add_shape(1, Inches(10.5), Inches(y), Inches(2.0), Inches(0.42))
        chip.fill.solid(); chip.fill.fore_color.rgb = color; chip.line.fill.background()
        cp = chip.text_frame.paragraphs[0]; cp.text = status
        cp.font.size = Pt(13); cp.font.bold = True; cp.font.color.rgb = WHITE
        cp.alignment = PP_ALIGN.CENTER
        y += 0.62
    return slide
# ---------- SLIDE 1: Title ----------
add_title_slide(
    "Airport Ground Operations\nManagement System",
    "Full-Stack Operations Platform with AI-Driven Predictions  |  Capstone Project Presentation"
)

# ---------- SLIDE 2: Agenda ----------
add_bullet_slide("Agenda", [
    "The Problem — why ground ops teams need this",
    "System Overview & Architecture",
    "Core Modules Walkthrough",
    "AI Module — predictions & chatbot",
    "Engineering Practices — branching, deployment, security",
    "Testing — security, performance, UI, UAT",
    "Results & Bugs Found",
    "Live Demo",
    "Conclusion & Next Steps",
])
# ---------- SLIDE 3: The Problem ----------
add_bullet_slide("The Problem", [
    "Airports coordinate flights, gates, maintenance, baggage, and staff across many teams in real time",
    "Poor coordination between these functions causes:",
    "  Flight delays and gate conflicts",
    "  Understaffed shifts",
    "  Slow response to equipment/maintenance issues",
    "These directly hurt on-time performance and passenger experience",
    "Most of this is still tracked over spreadsheets and radio calls",
], note="Source: docs/BUSINESS_OBJECTIVES.md")

# ---------- SLIDE 4: Measurable Business Goals ----------
add_status_slide("Measurable Business Goals", [
    ("Double-booked gate incidents / month", "TARGET: 0", SKY),
    ("Avg. time to detect turnaround delay", "< 5 MIN", SKY),
    ("Overdue maintenance requests", "-50%", SKY),
    ("Shift coverage gaps", "-30%", SKY),
    ("Incident log → assignment time", "< 10 MIN", SKY),
    ("Shifts planned using AI staffing prediction", "80%+", SKY),
])
# ---------- SLIDE 5: System Overview ----------
add_bullet_slide("System Overview", [
    "A full-stack web platform for day-to-day airport ground operations",
    "Covers: Flights, Gates, Staff, HR, Baggage, Maintenance, Ground Equipment, Notifications, Reports",
    "Integrated AI module for predictions + an operational chatbot",
    "Role-based access control, JWT auth, rate limiting, and full audit logging on every action",
    "REST API documented with Swagger/ReDoc",
])

# ---------- SLIDE 6: Tech Stack ----------
add_two_col_slide("Tech Stack",
    "Backend", [
        "Django 5 + Django REST Framework",
        "SimpleJWT (access + refresh tokens)",
        "PostgreSQL (prod) / SQLite (dev)",
        "django-ratelimit for login protection",
        "Swagger UI / ReDoc via drf-yasg",
    ],
    "Frontend & ML", [
        "React 19 + Vite",
        "Tailwind CSS 4, Recharts",
        "scikit-learn (RandomForest models)",
        "Rule-based chatbot (intent + live DB lookups)",
        "GitHub Actions CI/CD (tests + flake8)",
    ])
# ---------- SLIDE 7: Architecture ----------
slide = prs.slides.add_slide(BLANK)
_title_bar(slide, "Architecture")

def box_shape(x, y, w, h, text, fill, text_color=WHITE, size=15):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.color.rgb = WHITE
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size)
    p.font.bold = True; p.font.color.rgb = text_color; p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return s

box_shape(0.6, 1.8, 3.2, 1.1, "React + Vite SPA\n(Tailwind CSS)", SKY)
box_shape(5.0, 1.8, 3.2, 1.1, "Django + DRF API\n(JWT Auth)", NAVY)
box_shape(1.6, 3.6, 2.6, 1.0, "PostgreSQL /\nSQLite", NAVY)
box_shape(5.0, 3.6, 3.2, 1.0, "AI Module\n(scikit-learn + chatbot)", SKY)
box_shape(9.0, 3.6, 2.9, 1.0, "Audit Logging +\nNotifications", NAVY)

arrow_note = slide.shapes.add_textbox(Inches(3.85), Inches(2.0), Inches(1.1), Inches(0.7))
arrow_note.text_frame.paragraphs[0].text = "REST\nAPI ⇄"
arrow_note.text_frame.paragraphs[0].font.size = Pt(12)
arrow_note.text_frame.paragraphs[0].font.color.rgb = GREY

note = slide.shapes.add_textbox(Inches(0.8), Inches(5.1), Inches(11.5), Inches(1.5))
np_ = note.text_frame; np_.word_wrap = True
p = np_.paragraphs[0]
p.text = "Frontend is a separate SPA — talks to Django purely through the REST API, not server-rendered"
p.font.size = Pt(16); p.font.color.rgb = NAVY

# ---------- SLIDE 8: Core Modules ----------
add_two_col_slide("Core Modules (20+ Django Apps)",
    "Operations", [
        "accounts — auth, roles, JWT",
        "flights — airlines, aircraft, ground-handling pipeline",
        "gates — assignment, double-booking prevention",
        "staff — profiles, shifts, schedules",
        "baggage — check-in to claim, status history",
        "maintenance — requests & logs",
    ],
    "Supporting", [
        "hr_management — departments, leave",
        "ground_equipment — inventory & usage",
        "notifications — in-app alerts",
        "reports — operational reporting",
        "core_app — shared utils, permissions, audit log",
        "ai_module — ML predictions + chatbot",
    ])
# ---------- SLIDE 9: AI Module ----------
add_two_col_slide("AI Module — 7 ML Models + Chatbot",
    "Predictive Models (RandomForest)", [
        "Flight Delay — minutes + likelihood",
        "Predictive Maintenance — urgency score",
        "Passenger Rush — expected rush factor",
        "Weather Risk — risk score + delay flag",
        "Staff Requirement — crew/security/baggage counts",
        "Gate Recommendation — best-fit gate, ranked",
        "Equipment Failure — failure risk",
    ],
    "Chatbot", [
        "Rule-based: regex + intent matching",
        "Answers questions using live DB data",
        "Architected to be swappable for a real LLM later",
        "Delay & maintenance models auto-switch from",
        "synthetic → real historical data once enough",
        "records exist in the DB",
    ])

# ---------- SLIDE 10: Flight Status Pipeline ----------
slide = prs.slides.add_slide(BLANK)
_title_bar(slide, "Flight Ground-Handling Pipeline")
stages = ["Scheduled", "Gate\nAssigned", "Crew\nAssigned", "Fueling", "Boarding", "Departed"]
x = 0.6
for i, s in enumerate(stages):
    box_shape(x, 3.0, 1.85, 1.1, s, SKY if i % 2 == 0 else NAVY, size=14)
    if i < len(stages) - 1:
        arr = slide.shapes.add_textbox(Inches(x + 1.85), Inches(3.25), Inches(0.3), Inches(0.6))
        arr.text_frame.paragraphs[0].text = "→"
        arr.text_frame.paragraphs[0].font.size = Pt(22)
        arr.text_frame.paragraphs[0].font.color.rgb = GREY
    x += 2.15
sub = slide.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(11), Inches(0.8))
sub.text_frame.paragraphs[0].text = "Every transition is tracked, timestamped, and written to the audit log"
sub.text_frame.paragraphs[0].font.size = Pt(16)
sub.text_frame.paragraphs[0].font.color.rgb = NAVY
# ---------- SLIDE 11: Engineering Practices ----------
add_two_col_slide("Engineering Practices",
    "Branching Strategy (Git Flow)", [
        "main — production-ready only",
        "develop — integration branch",
        "feature/<desc> — one branch per task",
        "bugfix/ & hotfix/ branches",
        "PR + 1 reviewer approval required",
        "CI (GitHub Actions) must pass before merge",
    ],
    "Deployment", [
        "Documented deploy guide (docs/DEPLOYMENT.md)",
        "Frontend: static host (Netlify/Vercel) OR",
        "  served by Django via WhiteNoise",
        "Bug fixed: STATICFILES_DIRS was pointing at",
        "  raw frontend/ source instead of frontend/dist,",
        "  causing Django to try collecting React source",
        "  as static assets — corrected in settings.py",
    ])

# ---------- SLIDE 12: Section — Testing ----------
add_section_header("Testing & Quality Assurance")
# ---------- SLIDE 13: Security Testing ----------
add_two_col_slide("Security Testing",
    "Static Analysis & Dependencies", [
        "Bandit: 0 high, 0 medium severity",
        "  issues across 11,953 lines scanned",
        "55 low-severity informational-only findings",
        "pip-audit: 12 known CVEs found → 0 after",
        "  upgrading pip, requests, setuptools",
    ],
    "Application Security Suite", [
        "10/10 custom security tests passing",
        "Auth required on all list endpoints",
        "Invalid/malformed JWTs rejected",
        "Role-based write restrictions enforced",
        "Login rate-limited after 5 failed attempts",
        "SECRET_KEY never leaks in error responses",
    ])

# ---------- SLIDE 14: Performance Testing ----------
add_bullet_slide("Performance Testing (Locust)", [
    "Tool: Locust, 20 concurrent simulated users",
    "Clean run: 456 requests, 0 failures",
    "Median response time: 42ms  |  95th percentile: 370ms",
    "Login endpoint (/api/token/) slowest at ~1.2s median — expected",
    "  (deliberate password-hashing cost) — flagged for future monitoring",
    "🐛 Real bug found & fixed: RATELIMIT_ENABLE setting was ignoring",
    "  explicit environment variable overrides — now correctly configurable",
], note="Full detail: docs/PERFORMANCE_TEST_REPORT.md")
# ---------- SLIDE 15: Frontend UI Testing ----------
add_bullet_slide("Frontend UI Testing", [
    "Tool: Vitest + React Testing Library + user-event",
    "14/14 tests passing across Dashboard.jsx, Signup.jsx, Login.jsx",
    "Covers: loading states, data rendering, error handling,",
    "  empty states, form validation, API integration (mocked)",
    "Representative coverage — 3 of ~30+ page components tested,",
    "  documented as intentional scope in UI_TEST_REPORT.md",
    "🐛 Real bug found & fixed: form labels in Signup.jsx/Login.jsx",
    "  weren't linked to their inputs (htmlFor/id mismatch) —",
    "  an accessibility gap, fixed across 7 affected fields",
])

# ---------- SLIDE 16: Bugs Found Summary ----------
slide = prs.slides.add_slide(BLANK)
_title_bar(slide, "Real Bugs Found Across Testing Layers")
rows = [
    ("SECRET_KEY below recommended HMAC/SHA256 length", "Security Testing"),
    ("RATELIMIT_ENABLE ignored env var overrides", "Performance Testing"),
    ("Form labels not linked to inputs (accessibility)", "UI Testing"),
]
y = 1.8
for bug, source in rows:
    b1 = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(7.8), Inches(0.7))
    p = b1.text_frame.paragraphs[0]; p.text = "• " + bug
    p.font.size = Pt(18); p.font.color.rgb = NAVY
    b2 = slide.shapes.add_textbox(Inches(8.9), Inches(y), Inches(3.5), Inches(0.7))
    p2 = b2.text_frame.paragraphs[0]; p2.text = source
    p2.font.size = Pt(15); p2.font.italic = True; p2.font.color.rgb = SKY
    y += 0.95
note = slide.shapes.add_textbox(Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.3))
np_ = note.text_frame; np_.word_wrap = True
np_.paragraphs[0].text = "Three different testing layers caught three different classes of problems — a crypto/config issue, a settings bug, and an accessibility gap — none of which the others would have surfaced."
np_.paragraphs[0].font.size = Pt(15); np_.paragraphs[0].font.color.rgb = GREY
# ---------- SLIDE 17: UAT ----------
add_bullet_slide("User Acceptance Testing (UAT)", [
    "30+ manual test scenarios across real user workflows",
    "Roles covered: Admin, Operations Manager, Gate Manager,",
    "  HR, Ground Staff, Viewer (read-only)",
    "Covers: authentication, flight ground-ops workflow, gate",
    "  assignment, staff scheduling, and more",
    "Verifies real-world usability — distinct from automated",
    "  unit/integration/security/performance tests",
    "Includes a sign-off section for formal acceptance",
], note="Full checklist: docs/UAT_CHECKLIST.md")

# ---------- SLIDE 18: Overall Delivery Status ----------
add_status_slide("Delivery Status vs. Original Scope", [
    ("Business Objectives doc", "DONE", GREEN),
    ("Branching strategy / CONTRIBUTING.md", "DONE", GREEN),
    ("Deployment guide + WhiteNoise fix", "DONE", GREEN),
    ("Security testing (bandit + pip-audit + suite)", "DONE", GREEN),
    ("Performance testing (Locust)", "DONE", GREEN),
    ("Frontend UI testing", "DONE", GREEN),
    ("UAT checklist", "DONE", GREEN),
    ("Combined testing summary", "DONE", GREEN),
    ("Project presentation", "DONE", GREEN),
    ("Live demonstration", "AT PRESENTATION", SKY),
])
# ---------- SLIDE 19: Scope Notes (honest, not hidden) ----------
add_two_col_slide("Scope Notes — Documented, Not Hidden",
    "Services Layer", [
        "Business logic (multi-step validation,",
        "cross-model rules) moved out of views.py",
        "into a services.py pattern per app",
        "Refactored for 3 representative apps:",
        "  flights, gates, staff",
        "Pattern documented in ARCHITECTURE.md",
        "so it can be extended to remaining apps",
    ],
    "Frontend UI Test Coverage", [
        "14 tests across 3 of ~30+ page components:",
        "  Dashboard, Signup, Login",
        "Chosen as representative — auth flow +",
        "  a data-heavy page",
        "Documented as intentional scope in",
        "  UI_TEST_REPORT.md, not a gap discovered",
        "  after the fact",
    ])

# ---------- SLIDE 20: Conclusion ----------
add_bullet_slide("Conclusion & Next Steps", [
    "Delivered a full-stack ops platform: 20+ modules, JWT auth, audit",
    "  logging, 7 ML models, and a live chatbot",
    "Backed by real engineering practice: branching strategy, CI/CD,",
    "  documented deployment, and a services-layer architecture pattern",
    "Validated across 4 testing layers — security, performance, UI, UAT —",
    "  which together surfaced and fixed 3 real bugs",
    "Next steps: extend services layer + UI test coverage to remaining",
    "  apps, add frontend SAST, and expand rate limiting to AI endpoints",
])

# ---------- SLIDE 21: Thank You / Demo ----------
add_title_slide("Live Demo", "Thank you — questions welcome")
prs.save("Airport_Ground_Ops_Presentation.pptx")
print("Saved: Airport_Ground_Ops_Presentation.pptx")